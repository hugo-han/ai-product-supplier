#!/usr/bin/env python3
"""
Script: generate_pptx.py
Usage: pip install python-pptx cairosvg
Run: python scripts/generate_pptx.py --slides_dir docs/presentations --output docs/presentations/aibp-report.pptx

This script converts sample_slideN.svg to PNG and inserts them as full-slide images, and also adds speaker notes from slides_text.json.
"""
import os
import sys
import json
import argparse
from pptx import Presentation
from pptx.util import Inches
from cairosvg import svg2png


def svg_to_png(svg_path, png_path):
    svg2png(url=svg_path, write_to=png_path)


def build_pptx(slides_dir, output_path):
    json_path = os.path.join(slides_dir, 'slides_text.json')
    if not os.path.exists(json_path):
        print('slides_text.json not found in', slides_dir)
        return
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    prs = Presentation()
    # set 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for s in data['slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        svg_name = f"sample_slide{s['num']}.svg"
        svg_path = os.path.join(slides_dir, svg_name)
        if os.path.exists(svg_path):
            png_path = os.path.join(slides_dir, f"sample_slide{s['num']}.png")
            try:
                svg_to_png(svg_path, png_path)
                slide.shapes.add_picture(png_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            except Exception as e:
                print('Failed to convert and add', svg_path, e)
        # add speaker notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        bullets = '\n'.join(s.get('bullets', []))
        prompt = s.get('prompt', '')
        text_frame.text = f"{s.get('title','')}\n\nBullets:\n{bullets}\n\nPrompt:\n{prompt}"
    prs.save(output_path)
    print('Saved PPTX to', output_path)


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--slides_dir', default='docs/presentations')
    parser.add_argument('--output', default='docs/presentations/aibp-report.pptx')
    args = parser.parse_args()
    build_pptx(args.slides_dir, args.output)
